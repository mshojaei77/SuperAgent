import os
import hashlib
import fitz  # PyMuPDF
import pdfplumber
import docx
import docx2txt
from pptx import Presentation
from PIL import Image
import easyocr
import pandas as pd
from abc import ABC, abstractmethod
from settFILES_DIRings import 


# Abstract base class for all file readers
class IFileReader(ABC):
    @abstractmethod
    def read(self, file_path):
        pass


class TextFileReader(IFileReader):
    def read(self, file_path):
        with open(file_path, 'r') as file:
            content = file.read()
        return content


class PdfFileReader(IFileReader):
    def __init__(self, ocr_reader=None, extract_images=True, extract_text=True):
        self.ocr_reader = ocr_reader
        self.extract_images = extract_images
        self.extract_text = extract_text

    def read(self, file_path):
        text, images, image_text = "", [], ""
        image_dir = os.path.join(os.path.dirname(file_path), 'extracted_images')
        os.makedirs(image_dir, exist_ok=True)

        if self.extract_text:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
                    for table in page.extract_tables() or []:
                        text += "Table:\n" + "\n".join("\t".join(str(cell) for cell in row) for row in table) + "\n"

        if self.extract_images:
            pdf_document = fitz.open(file_path)
            for page in pdf_document:
                for img in page.get_images(full=True):
                    base_image = pdf_document.extract_image(img[0])
                    image_bytes = base_image["image"]
                    image_hash = hashlib.md5(image_bytes).hexdigest()
                    image_filename = os.path.join(image_dir, f"{image_hash}.{base_image['ext']}")
                    if not os.path.exists(image_filename):
                        with open(image_filename, "wb") as image_file:
                            image_file.write(image_bytes)
                        images.append(image_filename)

            if self.ocr_reader:
                for image_path in images:
                    image_text += f"\nText from {os.path.basename(image_path)}:\n"
                    image_text += ImageFileReader(self.ocr_reader).read(image_path)

        return {'text': text, 'images': images, 'image_text': image_text}


class DocxFileReader(IFileReader):
    def read(self, file_path):
        text = docx2txt.process(file_path)
        doc = docx.Document(file_path)
        for table in doc.tables:
            text += "\nTable:\n" + "\n".join("\t".join(cell.text for cell in row.cells) for row in table.rows) + "\n"
        return text


class PptxFileReader(IFileReader):
    def read(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text


class ImageFileReader(IFileReader):
    def __init__(self, ocr_reader=None):
        self.ocr_reader = ocr_reader

    def read(self, file_path):
        if self.ocr_reader:
            result = self.ocr_reader.readtext(file_path, detail=0)
            return "\n".join(result)
        return ""


class XlsxFileReader(IFileReader):
    def read(self, file_path):
        xl = pd.ExcelFile(file_path)
        text = ""
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text += f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n"
        return text


class FileManager:
    def __init__(self, ocr_languages=['en', 'fa'], enable_ocr=True):
        self.ocr_reader = easyocr.Reader(ocr_languages) if enable_ocr else None
        self.readers = {
            '.txt': TextFileReader(),
            '.pdf': PdfFileReader(self.ocr_reader),
            '.docx': DocxFileReader(),
            '.pptx': PptxFileReader(),
            '.xlsx': XlsxFileReader(),
            '.png': ImageFileReader(self.ocr_reader),
            '.jpg': ImageFileReader(self.ocr_reader),
            '.jpeg': ImageFileReader(self.ocr_reader),
            '.tiff': ImageFileReader(self.ocr_reader),
            '.bmp': ImageFileReader(self.ocr_reader),
            '.gif': ImageFileReader(self.ocr_reader),
        }

    def read(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext in self.readers:
            return self.readers[ext].read(file_path)
        else:
            print(f"Unsupported file type: {ext}")
            return None

    def cleanup(self, file_path):
        if os.path.exists(file_path):
            os.remove(file_path)
        image_dir = os.path.join(os.path.dirname(file_path), 'extracted_images')
        if os.path.exists(image_dir):
            for image_file in os.listdir(image_dir):
                image_file_path = os.path.join(image_dir, image_file)
                if os.path.isfile(image_file_path):
                    os.remove(image_file_path)
            os.rmdir(image_dir)


# Example usage:
if __name__ == "__main__":
    file_manager = FileManager(ocr_languages=['en', 'fa'], enable_ocr=False)
    
    file_uploaded = False
    
    directory_path = FILES_DIR
    
    filename = input("Please upload a file (example: sample.pdf): ").strip()
    
    if filename:
        file_uploaded = True
        
        
    if filename.lower().endswith(('.txt', '.pdf', '.docx', '.pptx', '.xlsx', '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif')):
        valid_file = True
    else:
        valid_file = False
    
    if file_uploaded & valid_file:
        file_path = os.path.join(directory_path, filename)
        
        content = file_manager.read(file_path)
        #file_manager.cleanup(file_path)
        print(content)
    else:
        print("Invalid file type.")
